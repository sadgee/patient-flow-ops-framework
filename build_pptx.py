"""Generate the team's final 15-slide PowerPoint deck.

Uses python-pptx with a clean modern light theme. Each slide has:
  - Title bar (color accent)
  - Body content (bullets, tables, or figures)
  - Speaker notes embedded in the notes pane (so the presenter can read them)

Run:
    .venv/bin/python analysis/build_pptx.py
Output: /Users/mac/Projects/Viharika/Team_A11_PatientFlow.pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "analysis" / "figures"
OUT = ROOT / "Team_A11_PatientFlow.pptx"

# Theme
COLOR_ACCENT = RGBColor(0x1F, 0x4E, 0x79)        # navy
COLOR_ACCENT_LIGHT = RGBColor(0xD9, 0xE4, 0xF1)
COLOR_TEXT = RGBColor(0x2B, 0x2B, 0x2B)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_HIGHLIGHT = RGBColor(0xE7, 0x6F, 0x51)
COLOR_GREEN = RGBColor(0x2A, 0x9D, 0x8F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, *, size=18, bold=False, color=COLOR_TEXT, name="Calibri"):
    run.text = text
    f = run.font
    f.name = name
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color


def add_accent_bar(slide, top=Inches(0.0), height=Inches(0.5), color=COLOR_ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_title(slide, title_text, subtitle_text=None):
    """Add a colored title bar with white title text."""
    add_accent_bar(slide, top=Inches(0.0), height=Inches(0.85))
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.10), SLIDE_W - Inches(1.0), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    set_run(p.add_run(), title_text, size=26, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle_text:
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), subtitle_text, size=14, color=RGBColor(0xDD, 0xDD, 0xDD))


def add_footer(slide, slide_num, total):
    tb = slide.shapes.add_textbox(Inches(0.4), SLIDE_H - Inches(0.4),
                                   SLIDE_W - Inches(0.8), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(),
            f"Team A11 — A Data-Driven Operations Framework for Patient Flow Management",
            size=10, color=COLOR_MUTED)
    p.alignment = PP_ALIGN.LEFT

    tb2 = slide.shapes.add_textbox(SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4),
                                    Inches(1.2), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    set_run(p2.add_run(), f"{slide_num} / {total}", size=10, color=COLOR_MUTED)
    p2.alignment = PP_ALIGN.RIGHT


def add_speaker_notes(slide, notes_text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes_text


def add_bullets(slide, bullets, *, left=Inches(0.6), top=Inches(1.2),
                  width=None, height=None, size=16):
    if width is None: width = SLIDE_W - Inches(1.2)
    if height is None: height = Inches(5.5)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            text, opts = b
        else:
            text, opts = b, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = opts.get("level", 0)
        bullet_marker = "▸ " if p.level == 0 else "•   "
        set_run(p.add_run(), bullet_marker + text,
                 size=opts.get("size", size),
                 bold=opts.get("bold", False),
                 color=opts.get("color", COLOR_TEXT))
        p.space_after = Pt(8)


def add_image(slide, image_path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(image_path), left, top, width=width)
    elif height:
        slide.shapes.add_picture(str(image_path), left, top, height=height)
    else:
        slide.shapes.add_picture(str(image_path), left, top)


def add_caption(slide, text, left, top, width):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    set_run(p.add_run(), text, size=10, color=COLOR_MUTED)
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, headers, rows, left, top, width, height,
                header_color=COLOR_ACCENT, header_text_color=RGBColor(0xFF, 0xFF, 0xFF),
                first_col_bold=False, highlight_rows=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    # headers
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.text = ""
        cell.text_frame.text = ""
        p = cell.text_frame.paragraphs[0]
        set_run(p.add_run(), h, size=12, bold=True, color=header_text_color)
        p.alignment = PP_ALIGN.CENTER
    # rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text_frame.text = ""
            p = cell.text_frame.paragraphs[0]
            bold = first_col_bold and j == 0
            color = COLOR_TEXT
            if highlight_rows and i - 1 in highlight_rows:
                color = COLOR_HIGHLIGHT
                bold = True
            set_run(p.add_run(), str(val), size=11, bold=bold, color=color)
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
    return tbl


# ----- Slide builders -----

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background accent
    add_accent_bar(s, top=Inches(0.0), height=SLIDE_H, color=COLOR_ACCENT)
    # White inner card
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.5),
                                SLIDE_W - Inches(2.0), SLIDE_H - Inches(3.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.fill.background()

    tb = s.shapes.add_textbox(Inches(1.4), Inches(2.2),
                                 SLIDE_W - Inches(2.8), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), "A Data-Driven Operations Framework",
              size=36, bold=True, color=COLOR_ACCENT)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), "for Patient Flow Management",
              size=36, bold=True, color=COLOR_ACCENT)
    p2.alignment = PP_ALIGN.CENTER

    tb2 = s.shapes.add_textbox(Inches(1.4), Inches(4.4),
                                  SLIDE_W - Inches(2.8), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    set_run(p3.add_run(), "Team A11", size=20, bold=True, color=COLOR_TEXT)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf2.add_paragraph()
    set_run(p4.add_run(),
              "Amitesh Mohan  ·  Gitanjali Roy  ·  Sadgee Pandey  ·  Viharika Appaneravanda",
              size=14, color=COLOR_TEXT)
    p4.alignment = PP_ALIGN.CENTER
    p5 = tf2.add_paragraph()
    set_run(p5.add_run(), "April 2026", size=12, color=COLOR_MUTED)
    p5.alignment = PP_ALIGN.CENTER

    add_speaker_notes(s, "Welcome the audience. Set up: outpatient clinics waste capacity even when full. Insurance constraints + multi-stage flow create non-obvious bottlenecks. Our framework quantifies the wait-vs-throughput cost of any operational policy before deploying it.")
    return s


def slide_problem(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "The Business Problem",
                "Outpatient clinics waste capacity even when full")
    add_bullets(s, [
        ("U.S. outpatient no-show rates: 15–30% — slots wasted, but patients still wait hours",
            {"size": 18}),
        ("Insurance networks restrict which provider a patient can see — demand pools at "
            "in-network providers, leaving others idle", {"size": 18}),
        ("Each visit is multi-stage: check-in → nurse intake → physician. "
            "Bottlenecks shift between stages and across the day", {"size": 18}),
        ("This is a system-level operations problem, not an appointment-grid scheduling tweak",
            {"size": 18, "bold": True}),
    ])
    add_footer(s, n, total)
    add_speaker_notes(s, "Open with the paradox: how can a clinic have idle providers AND long patient waits? Because capacity isn't fungible — an uninsured patient can't see a physician who only accepts private insurance, even if that physician is sitting idle. We're addressing this as a system-level operations problem.")
    return s


def slide_approach(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Approach",
                "Calibrated multi-stage simulator → test operational policies on it")

    steps = [
        ("1", "Predict no-show probability", "Gradient-boosted classifier on 110k Kaggle appointments"),
        ("2", "Calibrate service times", "Lognormal fits to NAMCS 2019 physician-time data"),
        ("3", "Build synthetic clinic", "Real provider mix from CMS · payer mix from MEPS"),
        ("4", "Run discrete-event simulation", "SimPy multi-stage flow with no-show stochasticity"),
        ("5", "Compare policies", "Throughput · wait · MD/nurse utilization · overtime"),
    ]
    top = Inches(1.4)
    for i, (num, title, desc) in enumerate(steps):
        y = top + Inches(0.95 * i)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = COLOR_ACCENT
        circ.line.fill.background()
        circ.text_frame.text = ""
        p = circ.text_frame.paragraphs[0]
        set_run(p.add_run(), num, size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        p.alignment = PP_ALIGN.CENTER
        # title + desc
        tb = s.shapes.add_textbox(Inches(1.6), y, SLIDE_W - Inches(2.4), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        set_run(p1.add_run(), title, size=18, bold=True, color=COLOR_ACCENT)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), desc, size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "Five steps. Each is grounded in a real public dataset for the U.S. context — so even though we simulate, the parameters are not made up. We end with policy comparisons measured on the operational metrics the proposal called out.")
    return s


def slide_data(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Data Sources — Hybrid & Honest",
                "Four public datasets · operational integration is synthesized")
    headers = ["Source", "What it provides", "Used for"]
    rows = [
        ["Kaggle Medical Appointments\n(110,527 rows)",
            "Appointments + no-show labels +\npatient covariates",
            "No-show probability model"],
        ["CDC NAMCS 2019\n(8,250 visits)",
            "Empirical U.S. physician-time TIMEMD",
            "Service-time distributions"],
        ["AHRQ MEPS HC-243\n(22,431 persons, 2022)",
            "U.S. payer mix (private/public/uninsured)",
            "Insurance-network shares"],
        ["CMS Doctors & Clinicians\n(786,227 after filter)",
            "Real U.S. clinicians + specialty + group",
            "Provider population"],
    ]
    add_table(s, headers, rows,
                Inches(0.5), Inches(1.2), SLIDE_W - Inches(1.0), Inches(3.5),
                first_col_bold=True)

    # Honesty box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.0),
                                SLIDE_W - Inches(1.0), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.margin_top = tb.margin_bottom = Inches(0.1)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Honesty:  ", size=13, bold=True, color=COLOR_HIGHLIGHT)
    set_run(p.add_run(),
              "No public dataset integrates insurance × multi-stage operations. "
              "We synthesize the operational layer, calibrated to those four sources. "
              "Kaggle is Brazilian; we use only its relative no-show patterns (lead time, "
              "age, prior history — cross-nationally robust per the literature) and rescale "
              "absolute probabilities to the U.S. base rate of 18%.",
              size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "Four real sources cover demand, service times, payer mix, and provider population. We synthesize the integrated operational view on top — which is what the proposal called out as the plan in point 3. Kaggle is Brazilian; we use only its relative no-show patterns and rescale absolute probabilities to a U.S. base rate.")
    return s


def slide_drivers(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "No-Show Drivers",
                "Lead time dominates · same-day 5% no-show, 30+ day 33%")
    add_image(s, FIG / "noshow_lead_time.png", Inches(0.4), Inches(1.2),
                width=Inches(6.2))
    add_image(s, FIG / "noshow_drivers.png", Inches(6.8), Inches(1.2),
                width=Inches(6.2))
    # Bottom strip with key facts
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.4),
                                SLIDE_W - Inches(1.0), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "AUC = 0.74 ", size=14, bold=True, color=COLOR_GREEN)
    set_run(p.add_run(),
              "(gradient-boosted, patient-level held-out test — no leakage; "
              "logit baseline 0.73)", size=13, color=COLOR_TEXT)
    p2 = tb.add_paragraph()
    set_run(p2.add_run(), "Watch the SMS confound:  ", size=13, bold=True, color=COLOR_HIGHLIGHT)
    set_run(p2.add_run(),
              "raw rates suggest SMS makes no-shows worse (27% with SMS vs 17% without). "
              "It's selection bias — clinics send SMS to high-risk patients. The model "
              "controls for it.", size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "Our no-show model gets AUC of 0.74, in the middle of the published range. We split by patient ID — no patient appears in both training and testing — so AUC isn't inflated by leakage. Lead time is the single biggest predictor. Watch out for the SMS confound.")
    return s


def slide_service_times(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Service-Time Calibration",
                "Lognormal fits to NAMCS 2019 by specialty category")
    add_image(s, FIG / "service_times_by_specialty.png",
                Inches(0.4), Inches(1.15), width=SLIDE_W - Inches(0.8))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.6),
                                SLIDE_W - Inches(1.0), Inches(1.4))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(),
              "Primary care: median 19 min  ·  Surgical: 19 min  ·  Medical specialty: 22 min  ·  "
              "Right-skewed, long tail to 90 min",
              size=13, color=COLOR_TEXT)
    p2 = tb.add_paragraph()
    set_run(p2.add_run(),
              "Check-in (~3 min) and nurse intake (~10 min) calibrated to OR literature "
              "(Cayirli & Veral 2003, Gupta & Denton 2008) — NAMCS doesn't measure them.",
              size=12, color=COLOR_MUTED)
    add_footer(s, n, total)
    add_speaker_notes(s, "Service times are right-skewed. We fit lognormal distributions per specialty category from NAMCS — around 20-25 minutes for an average visit, with a long right tail. For check-in and nurse intake, NAMCS doesn't measure them, so we calibrated using ratios from the operations research literature.")
    return s


def slide_clinic(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "The Synthetic Clinic",
                "20 providers × 10 specialties × 3 insurance networks · 2,000-patient demand pool")
    add_image(s, FIG / "clinic_overview.png",
                Inches(0.3), Inches(1.15), width=SLIDE_W - Inches(0.6))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.5),
                                SLIDE_W - Inches(1.0), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "All 20 providers in network A (private)  ·  ",
              size=13, color=COLOR_TEXT)
    set_run(p.add_run(), "14 in B (public)  ·  ", size=13, color=COLOR_TEXT)
    set_run(p.add_run(), "only 5 in C (Medicaid)", size=13, bold=True,
              color=COLOR_HIGHLIGHT)
    set_run(p.add_run(), "  ← the operational constraint", size=13, color=COLOR_TEXT)
    p2 = tb.add_paragraph()
    set_run(p2.add_run(),
              "Patients matched to providers on BOTH network and specialty "
              "(no pediatric→cardiologist mismatches).",
              size=12, color=COLOR_MUTED)
    add_footer(s, n, total)
    add_speaker_notes(s, "Here is the synthetic clinic. The interesting feature is the asymmetric network coverage: only 5 of 20 providers accept Medicaid, even though those patients are 7% of the population. This is the demand-capacity mismatch the proposal called out.")
    return s


def slide_baseline(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Baseline Simulation",
                "263 visits/day, 26-min median wait, 62% MD util — but nurse stage is binding")
    headers = ["Metric", "Mean", "95% CI", "Note"]
    rows = [
        ["Completed visits / day", "262.4", "[248, 273]", ""],
        ["Median wait to physician", "25.8 min", "[22, 29]", ""],
        ["90th-percentile wait", "42.0 min", "", ""],
        ["Mean MD utilization", "60.7%", "", "physicians have idle time"],
        ["Mean nurse utilization", "~93%", "", "saturated → bottleneck"],
        ["Nurse overtime", "~0 hrs/day", "", ""],
        ["No-show rate (rescaled)", "17.5%", "", "U.S. base anchor"],
    ]
    add_table(s, headers, rows,
                Inches(0.5), Inches(1.3), SLIDE_W - Inches(1.0), Inches(3.6),
                first_col_bold=True, highlight_rows={4})

    # Sanity check note
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.2),
                                SLIDE_W - Inches(1.0), Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Sanity check (simulator validation):  ", size=13, bold=True, color=COLOR_GREEN)
    set_run(p.add_run(),
              "with 12 nurses (vs. 6) and zero no-shows, median wait collapses from "
              "67 min → 29 min — confirming the simulator is correct AND that the nurse "
              "stage is what's saturated, not randomness.",
              size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "The baseline. Median wait is 26 minutes, decent for a multispecialty clinic. But MD utilization is only 62% — physicians are sitting idle. Why don't more patients get through? Look at nurse utilization: 93%. The nurse stage is saturated. We validated this with a sanity run — doubling nurse capacity collapses the wait.")
    return s


def slide_optimization(prs, n, total):
    """Frame the prescriptive question in formal optimization language.

    Decision variables → Objective → Constraints → Method.
    NAR and RBS are positioned as two corner solutions of this π-space.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Optimization Framing",
                "The prescriptive question: choose a policy π under binding constraints")

    # Four panels in 2x2 layout — decision vars / objective / constraints / method
    card_w = (SLIDE_W - Inches(1.6)) / 2
    card_h = Inches(2.45)
    cards = [
        # (col, row, title, title_color, body_runs)
        (0, 0, "DECISION VARIABLES", COLOR_ACCENT, [
            ("π = (p_overbook, p_buffer, buffer_min, {c_s})", True, 13, COLOR_ACCENT),
            ("\n• p_overbook ∈ [0,1] — threshold above which a high-risk slot gets a same-network standby",
              False, 11, COLOR_TEXT),
            ("\n• p_buffer  ∈ [0,1] — threshold above which the following slot is left empty",
              False, 11, COLOR_TEXT),
            ("\n• buffer_min ∈ {0, 15, 30} — minutes of protective slack",
              False, 11, COLOR_TEXT),
            ("\n• c_s ∈ [c_min, c_max] — slots/day for provider s (future extension)",
              False, 11, COLOR_MUTED),
        ]),
        (1, 0, "OBJECTIVE", COLOR_GREEN, [
            ("max  α · E[throughput]  −  β · E[median wait]  −  γ · E[nurse OT]",
              True, 13, COLOR_GREEN),
            ("\nequivalently:", False, 11, COLOR_MUTED),
            ("\n  max  E[throughput]\n  s.t.  E[median wait] ≤ T*", True, 12, COLOR_TEXT),
            ("          E[nurse OT] ≤ H*", True, 12, COLOR_TEXT),
            ("\nA clinic chooses (α, β, γ) or (T*, H*) based on its priorities "
              "(revenue pressure vs. patient experience vs. labor budget).",
              False, 10, COLOR_MUTED),
        ]),
        (0, 1, "CONSTRAINTS (always binding)", COLOR_HIGHLIGHT, [
            ("Hard — structural:", True, 12, COLOR_HIGHLIGHT),
            ("\n• Network match: patient p → provider s only if p.network ∈ s.networks",
              False, 11, COLOR_TEXT),
            ("\n• Specialty match: p.speccat_needed = s.speccat", False, 11, COLOR_TEXT),
            ("\n• Capacity: 1 MD/patient at a time · 2 check-in · 6 nurses · 8-hr day · 30-min slots",
              False, 11, COLOR_TEXT),
            ("\nSoft — penalty in objective:", True, 12, COLOR_HIGHLIGHT),
            ("\n• Nurse OT > 0 taxed at γ · (revealed in post-analysis)",
              False, 11, COLOR_TEXT),
        ]),
        (1, 1, "METHOD", COLOR_ACCENT, [
            ("Simulation-based policy evaluation.", True, 13, COLOR_ACCENT),
            ("\nThe problem is stochastic (no-shows, service-time variance), multi-stage "
              "(check-in → nurse → MD), and integer-valued → no closed-form optimum. "
              "Monte-Carlo simulation of the π-space is the right tool for this class.",
              False, 11, COLOR_TEXT),
            ("\nNAR and RBS are two corner solutions of π:", True, 11, COLOR_HIGHLIGHT),
            ("\n  NAR: p_overbook = 0.30, p_buffer = ∞  → max-throughput corner",
              False, 10, COLOR_TEXT),
            ("\n  RBS: p_overbook = ∞, p_buffer = 0.30, buffer = 30  → min-wait corner",
              False, 10, COLOR_TEXT),
        ]),
    ]
    for col, row, header, hdr_color, body in cards:
        x = Inches(0.5) + col * (card_w + Inches(0.3))
        y = Inches(1.05) + row * (card_h + Inches(0.2))
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
        box.line.color.rgb = hdr_color
        box.line.width = Pt(1.25)
        # Header tab
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.35))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = hdr_color
        hdr.line.fill.background()
        hdr.text_frame.text = ""
        hp = hdr.text_frame.paragraphs[0]
        set_run(hp.add_run(), header, size=11, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
        hp.alignment = PP_ALIGN.LEFT
        hdr.text_frame.margin_left = Inches(0.15)
        hdr.text_frame.margin_top = Inches(0.04)
        # Body text
        tb = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.4),
                                     card_w - Inches(0.3), card_h - Inches(0.45))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.word_wrap = True
        # First run replaces default paragraph; subsequent runs extend it
        p = tf.paragraphs[0]
        first = True
        for txt, bold, size, color in body:
            if txt.startswith("\n"):
                txt = txt.lstrip("\n")
                p = tf.add_paragraph()
                first = True
            set_run(p.add_run(), txt, size=size, bold=bold, color=color)

    # Problem-class strip (sits just above the page footer so they don't collide)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(6.35),
                                SLIDE_W - Inches(1.0), Inches(0.50))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT
    box.line.fill.background()
    tb = box.text_frame; tb.margin_left = tb.margin_right = Inches(0.2)
    tb.margin_top = Inches(0.05); tb.margin_bottom = Inches(0.05)
    p = tb.paragraphs[0]
    set_run(p.add_run(),
              "Class of problem:  constrained stochastic optimization under multi-stage "
              "queueing — evaluated via simulation on a U.S.-calibrated synthetic clinic.",
              size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    p.alignment = PP_ALIGN.CENTER

    add_footer(s, n, total)
    add_speaker_notes(s,
        "Before diving into the two policies, here is the problem in formal operations-analytics language. "
        "Decision variables: the policy parameters — which appointments to overbook, which to buffer, and by how much. "
        "Objective: maximize throughput minus weighted penalties for patient wait and nurse overtime — or equivalently, maximize throughput subject to service-level constraints. "
        "Constraints come in two flavors: hard structural constraints — insurance-network match, specialty match, physical capacity — that must always hold; and soft penalties, like overtime cost. "
        "The problem is stochastic, multi-stage, and integer-valued, so there is no closed form — simulation is the right tool. "
        "NAR and RBS are two corner solutions of the policy space: NAR maximizes the throughput corner, RBS maximizes the wait-minimization corner. "
        "The next slides evaluate both.")
    return s


def slide_wait_by_stage(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Where the Wait Actually Happens",
                "Decompose patient wait by stage · nurse-intake queue dominates under stress")
    add_image(s, FIG / "scenario_wait_by_stage.png",
                Inches(0.4), Inches(1.15), width=Inches(8.0))
    # right-side narrative
    tb = s.shapes.add_textbox(Inches(8.6), Inches(1.4),
                                 Inches(4.4), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Most diagnostic chart of the project",
              size=14, bold=True, color=COLOR_ACCENT)
    p.space_after = Pt(8)

    items = [
        ("Baseline", "15 min total · spread across stages", COLOR_TEXT),
        ("NAR (overbook)", "65 min total · 47 min in nurse queue alone", COLOR_HIGHLIGHT),
        ("RBS (buffer)", "9 min total · nurse stage breathes", COLOR_GREEN),
    ]
    for label, desc, color in items:
        p1 = tf.add_paragraph()
        set_run(p1.add_run(), label, size=13, bold=True, color=color)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), "  " + desc, size=12, color=COLOR_TEXT)
        p2.space_after = Pt(6)

    p3 = tf.add_paragraph()
    set_run(p3.add_run(),
              "Aggressive overbooking moves the bottleneck from MD to nurse stage — "
              "a stage-shifting failure invisible to single-server queueing analysis.",
              size=12, color=COLOR_MUTED)
    add_footer(s, n, total)
    add_speaker_notes(s, "This is the most diagnostic chart of the project. Each bar is the average wait per patient, broken down by which stage. In baseline, waits are spread. Under aggressive overbooking, the nurse-intake queue explodes — patients are stuck at the front desk because nurses can't get them through fast enough. RBS is the cleanest.")
    return s


def slide_nar(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Policy A — Network-Aware Reassignment (NAR)",
                "π: p_overbook = 0.30 · targets the max-throughput corner of the policy space")

    add_bullets(s, [
        ("Original no-shows → standby fills the slot (throughput gain)", {"size": 14}),
        ("Both show → both seen (denser MD queue)", {"size": 14}),
    ], top=Inches(1.2), height=Inches(0.9))

    # Result table
    headers = ["Metric", "Baseline", "NAR", "Δ"]
    rows = [
        ["Throughput (visits/day)", "262", "319", "+22%"],
        ["Median wait", "26 min", "77 min", "+200%"],
        ["MD utilization", "61%", "74%", "+13 pp"],
        ["Nurse overtime / day", "0 hr", "1.3 hr", "new cost"],
    ]
    add_table(s, headers, rows,
                Inches(0.5), Inches(2.4), Inches(7.5), Inches(2.6),
                first_col_bold=True, highlight_rows={1, 3})

    add_image(s, FIG / "scenario_overtime.png",
                Inches(8.3), Inches(2.4), width=Inches(4.7))

    # Insight box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.4),
                                SLIDE_W - Inches(1.0), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFC, 0xE5, 0xDD)
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Insight:  ", size=14, bold=True, color=COLOR_HIGHLIGHT)
    set_run(p.add_run(),
              "Naive overbooking hits the nurse stage, not the MD stage. Throughput is "
              "bought with overtime cost — real labor expense and burnout risk.",
              size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "Policy A is the textbook overbooking play. It works for throughput — 22% more visits — and pushes MD utilization to 74%. But median wait triples and nurses run an hour-plus of overtime per day. In a multi-stage system, optimizing one stage can break another.")
    return s


def slide_rbs(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Policy B — Risk-Buffered Scheduling (RBS)",
                "π: p_buffer = 0.30, buffer_min = 30 · targets the min-wait corner of the policy space")

    add_bullets(s, [
        ("Schedule density reduced on volatile appointments → fewer cascade delays", {"size": 14}),
        ("Trades throughput for predictability", {"size": 14}),
    ], top=Inches(1.2), height=Inches(0.9))

    headers = ["Metric", "Baseline", "RBS", "Δ"]
    rows = [
        ["Median wait", "26 min", "19 min", "−26%"],
        ["90th-percentile wait", "42 min", "26 min", "−38%"],
        ["Throughput", "262 visits", "225 visits", "−14%"],
        ["MD utilization", "61%", "53%", "−8 pp"],
        ["Nurse overtime", "0 hr", "0 hr", "no cost"],
    ]
    add_table(s, headers, rows,
                Inches(2.0), Inches(2.4), Inches(9.0), Inches(3.0),
                first_col_bold=True, highlight_rows={0, 1, 4})

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.7),
                                SLIDE_W - Inches(1.0), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xDF, 0xF2, 0xEE)
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Insight:  ", size=14, bold=True, color=COLOR_GREEN)
    set_run(p.add_run(),
              "Opposite tradeoff to NAR — better patient experience and zero overtime, "
              "but 14% fewer completed visits.",
              size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "Policy B does the opposite — instead of squeezing patients in, leave slack around volatile appointments. Median wait drops by a quarter, the 90th-percentile wait drops 38%, nurse stage breathes. Cost: 14% fewer visits and lower MD utilization. Clean wait-vs-volume tradeoff.")
    return s


def slide_equity(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Equity View — Network Constraint",
                "Wait is similar across payers today, but creates latent fragility")
    add_image(s, FIG / "scenario_per_payer.png",
                Inches(0.3), Inches(1.15), width=SLIDE_W - Inches(0.6))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.5),
                                SLIDE_W - Inches(1.0), Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Visit shares (62 / 30 / 8) ≈ population shares (58 / 34 / 7). ",
              size=13, color=COLOR_TEXT)
    set_run(p.add_run(), "Today, supply roughly matches demand for each payer. ",
              size=13, color=COLOR_TEXT)
    p2 = tb.add_paragraph()
    set_run(p2.add_run(), "Latent fragility:  ", size=13, bold=True, color=COLOR_HIGHLIGHT)
    set_run(p2.add_run(),
              "a 20% surge in uninsured demand has nowhere to go (only 5 providers); "
              "the same surge in private demand distributes across 20.",
              size=13, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "The equity view. Today, in our calibrated baseline, waits are similar across payer types because supply matches demand. But the structural asymmetry — only 5 of 20 providers accepting uninsured patients — means a demand surge has no relief valve. Stress-testing demand surges by payer is the obvious next step.")
    return s


def slide_tradeoff(prs, n, total):
    """Big tradeoff scatter, headline at the bottom."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Tradeoff Frontier",
                "Three policies sit on three corners of the throughput-vs-wait plane")
    add_image(s, FIG / "scenario_tradeoff.png",
                Inches(2.4), Inches(1.05), height=Inches(5.0))

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(6.2),
                                SLIDE_W - Inches(1.0), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2); tb.margin_top = Inches(0.05)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Read it as a Pareto plane:  ", size=13, bold=True, color=COLOR_ACCENT)
    set_run(p.add_run(),
              "NAR (top-right) buys throughput with wait. RBS (bottom-left) "
              "buys patient experience with throughput. Baseline (middle) is the status quo. "
              "The right policy depends on what the clinic optimizes for.",
              size=12, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "The tradeoff frontier — three policies, three corners. NAR pushes throughput up at the cost of waits. RBS pushes waits down at the cost of throughput. Baseline sits in the middle. The choice depends on what the clinic optimizes for. Next slide: the most surprising finding from sensitivity analysis.")
    return s


def slide_sensitivity(prs, n, total):
    """Full-width sensitivity panel + nurse-staffing headline."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Sensitivity — Where the Real Lever Is",
                "Vary one parameter at a time; baseline policy held fixed")
    add_image(s, FIG / "sensitivity_panel.png",
                Inches(0.3), Inches(1.1), width=SLIDE_W - Inches(0.6))

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.6),
                                SLIDE_W - Inches(1.0), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFC, 0xE5, 0xDD)
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.2); tb.margin_top = Inches(0.1)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Headline:  ", size=15, bold=True, color=COLOR_HIGHLIGHT)
    set_run(p.add_run(),
              "Adding 1 nurse (6 → 7) cuts median wait from 26 → 21 min, "
              "almost matching the RBS policy — and without sacrificing throughput. ",
              size=13, color=COLOR_TEXT)
    p2 = tb.add_paragraph()
    set_run(p2.add_run(),
              "Software policies have a ceiling. Staffing breaks through it.",
              size=14, bold=True, color=COLOR_ACCENT)
    add_footer(s, n, total)
    add_speaker_notes(s, "The single most actionable insight from the whole project. We swept three parameters one at a time. No-show rate (left): higher no-shows actually reduce wait by relieving the nurse stage. MD time (middle): waits are sensitive but throughput barely moves — confirms nurse stage is binding. Nurse staffing (right): one extra nurse cuts wait by 5 min — almost matching RBS without giving up throughput. Software policies have a ceiling; staffing breaks through it.")
    return s


def slide_real_world(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Real-World Anchors",
                "These dynamics are visible in current U.S. healthcare ops debates")
    items = [
        ("Kaiser Permanente nurse strikes (2022, 2023)",
            "thousands walked out citing chronic understaffing — "
            "exactly the binding constraint our simulator surfaces"),
        ("CMS Medicaid Managed Care Network Adequacy Final Rule (May 2024)",
            "codifies maximum acceptable wait-time and provider-density standards "
            "for Medicaid networks — directly addresses the asymmetry we modeled"),
        ("Mayo Clinic / Cleveland Clinic published throughput studies",
            "report no-show rates of 18–25% and explicit risk-stratified overbooking — "
            "same NAR-style policy, same wait-time tradeoffs"),
        ("AHRQ outpatient wait-time data",
            "median wait to see a primary-care MD: 26 days for new patients; in-clinic wait "
            "20-25 min — within range of our simulated baseline"),
    ]
    top = Inches(1.3)
    for i, (head, body) in enumerate(items):
        y = top + Inches(1.25 * i)
        # Marker
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05),
                                    Inches(0.35), Inches(0.35))
        circ.fill.solid(); circ.fill.fore_color.rgb = COLOR_ACCENT
        circ.line.fill.background()
        tb = s.shapes.add_textbox(Inches(1.15), y, SLIDE_W - Inches(1.5), Inches(1.1))
        tf = tb.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        set_run(p1.add_run(), head, size=14, bold=True, color=COLOR_ACCENT)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), body, size=12, color=COLOR_TEXT)
    add_footer(s, n, total)
    add_speaker_notes(s, "The dynamics we modeled are visible in current U.S. healthcare news. Kaiser nurse strikes were about exactly the binding constraint we surfaced. CMS just issued a network-adequacy rule for Medicaid. Mayo and Cleveland Clinic studies report the same no-show rates and use the same overbook policy. Our baseline wait — 26 minutes — matches AHRQ's published range.")
    return s


def slide_closing(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Limitations · Future Work · Take-away")
    # Two columns
    # Limitations
    tb1 = s.shapes.add_textbox(Inches(0.5), Inches(1.2),
                                  Inches(6.0), Inches(4.0))
    tf1 = tb1.text_frame; tf1.word_wrap = True
    p = tf1.paragraphs[0]
    set_run(p.add_run(), "Honest Limitations", size=18, bold=True, color=COLOR_HIGHLIGHT)
    p.space_after = Pt(8)
    for txt in [
        "Demand calibrated on Brazilian Kaggle; absolute level rescaled to U.S.",
        "Operational layer (network labels, multi-stage times) synthesized — "
            "real EHR + insurance contracts would refine it",
        "Single-day simulation; no demand carryover or no-show callback dynamics",
        "No financial layer (revenue per visit, OT cost) — easy to add given a clinic's contracts",
    ]:
        p1 = tf1.add_paragraph()
        set_run(p1.add_run(), "▸ " + txt, size=12, color=COLOR_TEXT)
        p1.space_after = Pt(4)

    # Future work
    tb2 = s.shapes.add_textbox(Inches(6.8), Inches(1.2),
                                  Inches(6.0), Inches(4.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]
    set_run(p.add_run(), "Future Work (priority order)", size=18, bold=True,
              color=COLOR_GREEN)
    p.space_after = Pt(8)
    for txt in [
        "Stress-test demand surges by payer — surface the equity-fragility risk",
        "Blended policy: overbook only above p > 0.5, buffer in 0.3–0.5 band",
        "Cross-trained nurse pool as third operational lever",
        "Multi-day demand carryover with no-show callback queues",
    ]:
        p1 = tf2.add_paragraph()
        set_run(p1.add_run(), "▸ " + txt, size=12, color=COLOR_TEXT)
        p1.space_after = Pt(4)

    # Take-away box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.4),
                                SLIDE_W - Inches(1.0), Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT
    box.line.fill.background()
    tb = box.text_frame
    tb.margin_left = tb.margin_right = Inches(0.25)
    tb.margin_top = Inches(0.15)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    set_run(p.add_run(), "Take-away  ", size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    set_run(p.add_run(),
              "A modular, data-grounded simulator can identify which stage actually "
              "constrains a clinic and quantify the wait-vs-throughput cost of any policy "
              "before deployment. ",
              size=13, color=RGBColor(0xFF, 0xFF, 0xFF))
    p2 = tb.add_paragraph()
    set_run(p2.add_run(),
              "The most actionable insight from our analysis: the binding constraint is "
              "rarely where the proposed intervention targets — and the simulator "
              "surfaces this in 30 seconds.",
              size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    add_footer(s, n, total)
    add_speaker_notes(s, "Closing. We're explicit about limitations — Brazilian no-show data rescaled, synthesized operational layer, single-day simulation, no financial layer. The real value is the framework: any clinic that gets us their EHR data can plug it in and get policy-specific quantified tradeoffs in days, not months. The single most actionable take-away from our work is that the binding constraint is rarely where the proposed intervention targets. Thanks — happy to take questions.")
    return s


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("title", slide_title),
        ("problem", slide_problem),
        ("approach", slide_approach),
        ("data", slide_data),
        ("drivers", slide_drivers),
        ("service_times", slide_service_times),
        ("clinic", slide_clinic),
        ("baseline", slide_baseline),
        ("optimization", slide_optimization),
        ("wait_by_stage", slide_wait_by_stage),
        ("nar", slide_nar),
        ("rbs", slide_rbs),
        ("equity", slide_equity),
        ("tradeoff", slide_tradeoff),
        ("sensitivity", slide_sensitivity),
        ("real_world", slide_real_world),
        ("closing", slide_closing),
    ]
    total = len(builders)
    for i, (name, fn) in enumerate(builders, start=1):
        if name == "title":
            fn(prs)
        else:
            fn(prs, i, total)

    prs.save(str(OUT))
    print(f"Wrote {OUT}  ({total} slides)")


if __name__ == "__main__":
    main()
